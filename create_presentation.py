import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_presentation():
    prs = pptx.Presentation()
    # Set 16:9 widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]
    
    # Colors
    NAVY = RGBColor(0x0F, 0x17, 0x2A)
    SLATE = RGBColor(0x1E, 0x29, 0x3B)
    BLUE = RGBColor(0x25, 0x63, 0xEB)
    GREEN = RGBColor(0x10, 0xB9, 0x81)
    LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
    GRAY = RGBColor(0x64, 0x74, 0x8B)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    
    def add_header(slide, title_text, category_text="RESTAURANT GROWTH ANALYTICS"):
        # Header banner shape
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY
        shape.line.color.rgb = NAVY
        
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11.7), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p_cat = tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = GREEN
        
        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE
        
    def add_card(slide, left, top, width, height, bg_color=LIGHT_BG, border_color=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.color.rgb = bg_color
        return shape

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = NAVY
    bg1.line.color.rgb = NAVY
    
    tbox1 = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf1 = tbox1.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "RESTAURANT GROWTH POTENTIAL MODELING & STRATEGIC CLASSIFICATION SYSTEM"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.alignment = PP_ALIGN.LEFT
    
    p2 = tf1.add_paragraph()
    p2.text = "A Data-Driven Decision Support System for Hospitality Expansion & Portfolio Management"
    p2.font.size = Pt(18)
    p2.font.color.rgb = GREEN
    p2.space_before = Pt(14)
    
    p3 = tf1.add_paragraph()
    p3.text = "MSc Statistics & Data Science Capstone Presentation | Dataset: SkyCity Auckland (N = 1,696)"
    p3.font.size = Pt(13)
    p3.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
    p3.space_before = Pt(24)

    # ==========================================
    # SLIDE 2: Introduction & Background
    # ==========================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide2, "Introduction & Business Context")
    
    c1 = add_card(slide2, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    tb = slide2.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Multi-Channel Food Service Landscape"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    bullets1 = [
        "Modern restaurants operate in complex environments blending physical dining, self-delivery fleets, and third-party delivery aggregators (UberEats, DoorDash).",
        "Traditional store evaluation relies heavily on gross revenue or accounting profit, ignoring channel-specific commission friction and unit economics.",
        "A Auckland-wide study across 1,696 outlets provides empirical ground truth for data-driven strategic expansion."
    ]
    for b in bullets1:
        bp = tf.add_paragraph()
        bp.text = f"• {b}"
        bp.font.size = Pt(13)
        bp.font.color.rgb = SLATE
        bp.space_before = Pt(10)
        
    c2 = add_card(slide2, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    tb2 = slide2.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.9))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "Strategic Analytics Solution"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    bullets2 = [
        "Holistic Growth Score (0–100): Combines top-line growth, net margin, revenue scale, ticket size, and cost rates into one explainable index.",
        "Unsupervised Clustering (K-Means): Identifies 5 distinct operational archetypes across the portfolio.",
        "Supervised ML Classifiers: Machine learning algorithms predicting strategic expansion readiness with 93.2% accuracy.",
        "Interactive Dashboard: Deploys actionable decision tools for executives and operators."
    ]
    for b in bullets2:
        bp = tf2.add_paragraph()
        bp.text = f"• {b}"
        bp.font.size = Pt(13)
        bp.font.color.rgb = SLATE
        bp.space_before = Pt(10)

    # ==========================================
    # SLIDE 3: Problem Statement
    # ==========================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide3, "Problem Statement: Key Analytical Deficiencies")
    
    probs = [
        ("1. Isolated Metric Evaluation", "Decision-makers evaluate revenue or growth rate in isolation, risking premature expansion of high-volume but loss-making stores."),
        ("2. Subjective Classification", "Absence of statistically standardized restaurant grading leads to arbitrary capital allocation and inconsistent store management."),
        ("3. Hidden Commission Drag", "Operators underestimate third-party aggregator commissions (27–33%), inflating perceived profitability of delivery revenue."),
        ("4. Generic Operational Advice", "Franchise holders receive static recommendations rather than automated, data-driven prescriptive strategies tailored to actual metrics.")
    ]
    coords = [(0.8, 1.6), (6.8, 1.6), (0.8, 4.4), (6.8, 4.4)]
    for idx, (title, desc) in enumerate(probs):
        l, t = coords[idx]
        add_card(slide3, Inches(l), Inches(t), Inches(5.7), Inches(2.5))
        tb = slide3.shapes.add_textbox(Inches(l+0.2), Inches(t+0.2), Inches(5.3), Inches(2.1))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLUE
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = SLATE
        p2.space_before = Pt(8)

    # ==========================================
    # SLIDE 4: Objectives
    # ==========================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide4, "Project Objectives & Scope")
    
    objs = [
        "1. Complete Data Hygiene: Clean and validate 1,696 restaurant records across 30 operational/financial attributes.",
        "2. Feature Engineering: Derive total revenue, net profit, margin %, delivery order share, and cost efficiency ratios.",
        "3. Growth Scoring (0–100): Develop explainable Multi-Criteria Decision Analysis (MCDA) scoring methodology.",
        "4. Operational Clustering: Perform K-Means clustering ($k=5$) with Silhouette optimization and 2D/3D PCA visualization.",
        "5. Strategic Classification: Formulate a 5-class strategic expansion taxonomy mapped by machine learning algorithms.",
        "6. Feature Contribution: Quantify primary growth drivers using Pearson/Spearman correlation and Gini importances.",
        "7. Supervised ML Modeling: Benchmark Logistic Regression, Gradient Boosting, Random Forest, and Decision Trees.",
        "8. Interactive Web Dashboard: Build a 7-tab Streamlit web application for executive decision support."
    ]
    
    add_card(slide4, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3))
    tb = slide4.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(11.1), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, o in enumerate(objs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = o
        p.font.size = Pt(13)
        p.font.color.rgb = SLATE
        if i > 0:
            p.space_before = Pt(6)

    # ==========================================
    # SLIDE 5: Dataset Description
    # ==========================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide5, "Empirical Dataset Overview (SkyCity Auckland)")
    
    metrics_summary = [
        ("Total Records", "1,696 Outlets"),
        ("Original Attributes", "30 Variables"),
        ("Avg Monthly Revenue", "$45,836.86"),
        ("Avg Monthly Net Profit", "$4,638.15"),
        ("Avg Profit Margin", "10.42%"),
        ("Avg Delivery Order Share", "82.48%")
    ]
    for idx, (lbl, val) in enumerate(metrics_summary):
        r_i = idx // 3
        c_i = idx % 3
        l = 0.8 + c_i * 3.95
        t = 1.5 + r_i * 1.5
        add_card(slide5, Inches(l), Inches(t), Inches(3.7), Inches(1.3))
        tb = slide5.shapes.add_textbox(Inches(l+0.1), Inches(t+0.1), Inches(3.5), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = lbl.upper()
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = GRAY
        p2 = tf.add_paragraph()
        p2.text = val
        p2.font.size = Pt(20)
        p2.font.bold = True
        p2.font.color.rgb = BLUE
        
    add_card(slide5, Inches(0.8), Inches(4.7), Inches(11.7), Inches(2.2))
    tb = slide5.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(11.3), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Dataset Composition & Diversity"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    d_bullets = [
        "Subregions (4): CBD, North Shore, South Auckland, West Auckland.",
        "Cuisine Specializations (8): Burgers, Chicken Dishes, Chinese, Indian, Japanese, Kebabs/Mediterranean, Pizza, Thai.",
        "Operational Segments (4): Cafe, QSR, Ghost Kitchen, Full-service.",
        "Order & Channel Coverage: In-store, UberEats, DoorDash, and Self-delivery revenue, orders, and net profits."
    ]
    for b in d_bullets:
        bp = tf.add_paragraph()
        bp.text = f"• {b}"
        bp.font.size = Pt(12)
        bp.font.color.rgb = SLATE
        bp.space_before = Pt(4)

    # ==========================================
    # SLIDE 6: Methodology Workflow
    # ==========================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide6, "Methodology & Analytical Pipeline")
    
    steps = [
        ("1. Data Preprocessing", "Missing check, duplicate removal, datatype casting, MinMax & Standard scaling."),
        ("2. Feature Engineering", "Derive total revenue, net profit, margin %, delivery share, & cost efficiency."),
        ("3. Growth Scoring", "Explainable 0–100 composite index via Multi-Criteria Decision Analysis."),
        ("4. Clustering (K-Means)", "Unsupervised grouping ($k=5$) validated by Silhouette score ($0.375$) & PCA."),
        ("5. Strategic Classification", "5-category taxonomy mapped via supervised ML classifiers (93.2% accuracy)."),
        ("6. Streamlit Deployment", "Interactive 7-tab executive web application with scenario recommendation engine.")
    ]
    for idx, (stitle, sdesc) in enumerate(steps):
        r_i = idx // 2
        c_i = idx % 2
        l = 0.8 + c_i * 5.95
        t = 1.5 + r_i * 1.8
        add_card(slide6, Inches(l), Inches(t), Inches(5.7), Inches(1.6))
        tb = slide6.shapes.add_textbox(Inches(l+0.2), Inches(t+0.15), Inches(5.3), Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = stitle
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = GREEN
        p2 = tf.add_paragraph()
        p2.text = sdesc
        p2.font.size = Pt(12)
        p2.font.color.rgb = SLATE
        p2.space_before = Pt(4)

    # ==========================================
    # SLIDE 7: Data Preprocessing & Feature Engineering
    # ==========================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide7, "Feature Engineering & Key Metrics")
    
    add_card(slide7, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3))
    tb = slide7.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(11.1), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Mathematical Definitions of Derived Features"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    fe_formulas = [
        "Total Revenue = Revenue(InStore) + Revenue(UberEats) + Revenue(DoorDash) + Revenue(SelfDelivery)",
        "Total Net Profit = Profit(InStore) + Profit(UberEats) + Profit(DoorDash) + Profit(SelfDelivery)",
        "Profit Margin (%) = [ Total Net Profit / Total Revenue ] × 100%",
        "Delivery Order Share = [ Delivery Orders / Total Monthly Orders ] (Mean = 82.48%)",
        "Cost Efficiency Ratio = 1.0 - ( COGS Rate + OPEX Rate ) (Mean = 31.0%)",
        "Third-Party Share = [ (UberEats Revenue + DoorDash Revenue) / Total Revenue ]",
        "Growth Rate (%) = [ Growth Factor - 1.0 ] × 100%"
    ]
    for f in fe_formulas:
        fp = tf.add_paragraph()
        fp.text = f"• {f}"
        fp.font.size = Pt(13)
        fp.font.color.rgb = SLATE
        fp.space_before = Pt(8)

    # ==========================================
    # SLIDE 8: Growth Potential Score
    # ==========================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide8, "Explainable 0–100 Growth Potential Score")
    
    add_card(slide8, Inches(0.8), Inches(1.5), Inches(6.5), Inches(5.3))
    tb = slide8.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(6.1), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Multi-Criteria Decision Weighting"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    weights = [
        "Growth Velocity Factor: 20%",
        "Total Net Profit ($): 20%",
        "Net Profit Margin (%): 15%",
        "Total Monthly Revenue ($): 15%",
        "Average Order Value (AOV): 10%",
        "Cost Efficiency (1 - COGS - OPEX): 10%",
        "Delivery Margin & Logistics: 10%"
    ]
    for w in weights:
        wp = tf.add_paragraph()
        wp.text = f"• {w}"
        wp.font.size = Pt(13)
        wp.font.color.rgb = SLATE
        wp.space_before = Pt(6)
        
    if os.path.exists("report_charts/growth_score_dist.png"):
        slide8.shapes.add_picture("report_charts/growth_score_dist.png", Inches(7.5), Inches(1.8), width=Inches(5.0))

    # ==========================================
    # SLIDE 9: Clustering Analysis
    # ==========================================
    slide9 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide9, "Operational Clustering & PCA Projections (k=5)")
    
    add_card(slide9, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.3))
    tb = slide9.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.4), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "5 Operational Archetypes (Silhouette = 0.375)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    c_list = [
        "1. Delivery Driven Powerhouses (N = 570, 33.6%): High delivery share (84.2%), avg rev $35,000.",
        "2. Expansion Leaders (N = 400, 23.6%): Top profit margin (18.3%), avg rev $61,276, profit $11,221.",
        "3. High Rev / Low Margin (N = 333, 19.6%): High rev ($51,605) but negative profit (-$2,904).",
        "4. Stable Performers (N = 274, 16.2%): Consistent profit ($6,259), solid cash generation.",
        "5. At-Risk / Low Margin (N = 119, 7.0%): Unprofitable (-$2,814), severe OPEX drag."
    ]
    for cl in c_list:
        cp = tf.add_paragraph()
        cp.text = cl
        cp.font.size = Pt(11.5)
        cp.font.color.rgb = SLATE
        cp.space_before = Pt(6)
        
    if os.path.exists("report_charts/cluster_pca_scatter.png"):
        slide9.shapes.add_picture("report_charts/cluster_pca_scatter.png", Inches(6.8), Inches(1.8), width=Inches(5.7))

    # ==========================================
    # SLIDE 10: Strategic Classification
    # ==========================================
    slide10 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide10, "Strategic Classification Taxonomy")
    
    cats_slide = [
        ("EXPAND AGGRESSIVELY (6.6%)", "Growth Score ≥ 70, Margin ≥ 10%, Profit > $5,000. Prime for franchise expansion.", GREEN),
        ("INVEST SELECTIVELY (32.2%)", "Growth Score 55–70, Margin ≥ 7.5%. Solid unit economics; targeted marketing investment.", BLUE),
        ("OPTIMIZE BEFORE EXPANSION (19.5%)", "Revenue ≥ $40k but Margin < 7.5%. High volume constrained by COGS/OPEX rates.", RGBColor(0xF5, 0x9E, 0x0B)),
        ("MAINTAIN (22.7%)", "Growth Score 42–55. Stable store operations providing steady cash flow.", GRAY),
        ("RESTRUCTURE / HIGH RISK (19.0%)", "Growth Score < 42 or negative profit. Requires menu & cost overhaul.", RGBColor(0xEF, 0x44, 0x44))
    ]
    for idx, (ctitle, cdesc, ccol) in enumerate(cats_slide):
        t = 1.5 + idx * 1.05
        add_card(slide10, Inches(0.8), Inches(t), Inches(6.2), Inches(0.95))
        tb = slide10.shapes.add_textbox(Inches(1.0), Inches(t+0.05), Inches(5.8), Inches(0.85))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ctitle
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = ccol
        p2 = tf.add_paragraph()
        p2.text = cdesc
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = SLATE
        
    if os.path.exists("report_charts/strategic_category_dist.png"):
        slide10.shapes.add_picture("report_charts/strategic_category_dist.png", Inches(7.3), Inches(1.8), width=Inches(5.2))

    # ==========================================
    # SLIDE 11: Feature Contribution / Growth Drivers
    # ==========================================
    slide11 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide11, "Feature Contribution & Growth Drivers")
    
    add_card(slide11, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.3))
    tb = slide11.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.1), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Primary Expansion Drivers"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    drivers_list = [
        "1. Growth Velocity Factor (16.16%): Strongest predictor of long-term store scaling.",
        "2. Total Net Profit (15.93%): Absolute cash generation capacity.",
        "3. Total Revenue Scale (14.11%): Top-line market order volume.",
        "4. Net Profit Margin % (12.07%): Unit economic efficiency.",
        "5. Pearson Correlation: Net Profit (r = 0.7859) and Profit Margin (r = 0.7365) exhibit highest positive correlation with Growth Score."
    ]
    for d in drivers_list:
        dp = tf.add_paragraph()
        dp.text = f"• {d}"
        dp.font.size = Pt(12.5)
        dp.font.color.rgb = SLATE
        dp.space_before = Pt(8)
        
    if os.path.exists("report_charts/feature_importance.png"):
        slide11.shapes.add_picture("report_charts/feature_importance.png", Inches(6.6), Inches(1.8), width=Inches(5.9))

    # ==========================================
    # SLIDE 12: Streamlit Dashboard
    # ==========================================
    slide12 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide12, "Interactive Streamlit Web Dashboard Architecture")
    
    dash_features = [
        ("Tab 1: Executive Overview", "Portfolio KPIs, strategic category donut charts, & subregion comparisons."),
        ("Tab 2: Cluster Map & PCA Explorer", "Interactive 2D/3D PCA scatter plots & Silhouette evaluation curves."),
        ("Tab 3: Growth Scorecards", "Single-restaurant 0–100 gauge charts, metrics, & radar charts vs. cuisine average."),
        ("Tab 4: Feature Contribution", "Random Forest Gini feature importances & Pearson correlation heatmaps."),
        ("Tab 5: Strategy Recommendations", "Automated prescriptive advice generator for individual restaurant outlets."),
        ("Tab 6: Restaurant Comparison", "Side-by-side metric comparison matrix across 2 to 5 selected outlets."),
        ("Tab 7: ML Model Evaluation", "Machine learning performance leaderboard & confusion matrix heatmap visualizer.")
    ]
    for idx, (dtitle, ddesc) in enumerate(dash_features):
        r_i = idx // 2
        c_i = idx % 2
        l = 0.8 + c_i * 5.95
        t = 1.5 + r_i * 1.3
        add_card(slide12, Inches(l), Inches(t), Inches(5.7), Inches(1.15))
        tb = slide12.shapes.add_textbox(Inches(l+0.15), Inches(t+0.1), Inches(5.4), Inches(0.95))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = dtitle
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = BLUE
        p2 = tf.add_paragraph()
        p2.text = ddesc
        p2.font.size = Pt(11)
        p2.font.color.rgb = SLATE

    # ==========================================
    # SLIDE 13: Key Results & ML Benchmarks
    # ==========================================
    slide13 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide13, "Key Results & Machine Learning Benchmark")
    
    add_card(slide13, Inches(0.8), Inches(1.5), Inches(6.0), Inches(5.3))
    tb = slide13.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.6), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Supervised Model Benchmark (80/20 Split)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    ml_rows = [
        "1. Logistic Regression: 93.24% Accuracy | 0.9324 F1 | 0.9952 ROC-AUC (Top Model)",
        "2. Gradient Boosting: 92.65% Accuracy | 0.9257 F1 | 0.9932 ROC-AUC",
        "3. Random Forest: 87.94% Accuracy | 0.8759 F1 | 0.9837 ROC-AUC",
        "4. Decision Tree: 86.76% Accuracy | 0.8659 F1 | 0.9448 ROC-AUC"
    ]
    for mr in ml_rows:
        mp = tf.add_paragraph()
        mp.text = f"• {mr}"
        mp.font.size = Pt(12)
        mp.font.color.rgb = SLATE
        mp.space_before = Pt(8)
        
    if os.path.exists("report_charts/confusion_matrix.png"):
        slide13.shapes.add_picture("report_charts/confusion_matrix.png", Inches(7.1), Inches(1.8), width=Inches(5.4))

    # ==========================================
    # SLIDE 14: Recommendations
    # ==========================================
    slide14 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide14, "Prescriptive Strategic Recommendations")
    
    recs_ppt = [
        ("Expand Aggressively Outlets", "Capitalize on strong brand equity; open new subregion sites, boost digital marketing by 15-20%, & standardize SOPs for franchising.", GREEN),
        ("High Revenue / Low Margin Outlets", "Freeze physical expansion. Renegotiate vendor supplier pricing to reduce COGS by 3-5%, audit labor schedules, & apply 10-15% delivery menu markups.", RGBColor(0xF5, 0x9E, 0x0B)),
        ("Delivery Channel Friction", "Promote direct self-delivery ordering to avoid third-party commission fees (27-33%) and improve net margin by 15-22%.", BLUE),
        ("At-Risk / Unprofitable Outlets", "Conduct immediate menu audit, eliminate loss-making items, trim OPEX overhead, & shift focus to direct in-store dining.", RGBColor(0xEF, 0x44, 0x44))
    ]
    for idx, (rtitle, rdesc, rcol) in enumerate(recs_ppt):
        t = 1.5 + idx * 1.3
        add_card(slide14, Inches(0.8), Inches(t), Inches(11.7), Inches(1.15))
        tb = slide14.shapes.add_textbox(Inches(1.1), Inches(t+0.1), Inches(11.1), Inches(0.95))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = rtitle
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = rcol
        p2 = tf.add_paragraph()
        p2.text = rdesc
        p2.font.size = Pt(12)
        p2.font.color.rgb = SLATE

    # ==========================================
    # SLIDE 15: Conclusion & Future Scope
    # ==========================================
    slide15 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide15, "Conclusion & Future Scope")
    
    add_card(slide15, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    tb = slide15.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Project Summary & Impact"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    concl_bullets = [
        "Successfully developed an end-to-end Data Science system replacing subjective store grading with statistical rigor.",
        "Integrated Growth Scoring, K-Means clustering ($k=5$), and Supervised ML classifiers (93.24% accuracy).",
        "Streamlit dashboard provides an intuitive executive decision-support platform for hospitality expansion."
    ]
    for b in concl_bullets:
        bp = tf.add_paragraph()
        bp.text = f"• {b}"
        bp.font.size = Pt(13)
        bp.font.color.rgb = SLATE
        bp.space_before = Pt(10)
        
    add_card(slide15, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    tb2 = slide15.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.9))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "Future Scope & Enhancements"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    fut_bullets = [
        "Time-Series Forecasting: Incorporate multi-period longitudinal data for ARIMA/LSTM revenue forecasting.",
        "Geospatial GIS Intelligence: Integrate exact GPS coordinates and foot-traffic data for site selection.",
        "Real-Time POS Connectors: Build direct API integrations for continuous automated store scoring.",
        "Cloud Microservices: Deploy containerized model endpoints for enterprise scalability."
    ]
    for b in fut_bullets:
        bp = tf2.add_paragraph()
        bp.text = f"• {b}"
        bp.font.size = Pt(13)
        bp.font.color.rgb = SLATE
        bp.space_before = Pt(10)
        
    ppt_path = "Restaurant_Growth_Potential_Presentation.pptx"
    prs.save(ppt_path)
    print(f"Presentation saved successfully as {ppt_path}!")

if __name__ == "__main__":
    create_presentation()
